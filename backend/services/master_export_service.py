"""
套用母版匯出（Apply-master export）

把每頁的「可編輯文字」與「AI 生成的內容圖形」放到使用者提供的 PowerPoint 母版上：
  - 底圖 / logo / 頁尾 / 字型佈景  → 來自母版（用母版的版面配置，保留其背景）
  - AI 圖的整頁背景 / logo / 烤死文字 → 丟棄（不放 AI 整張圖）
  - 內容圖形（插畫/icon/圖表）      → 用視覺模型偵測 + 去背後貼上
  - 標題與內文                       → 依 OCR bbox 放成可編輯文字框

視覺模型透過 llm-proxy 呼叫（Claude 走 Anthropic 視覺）。
"""
import io
import logging
import os
from typing import List, Dict, Optional

from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from utils.graphic_extractor import detect_logos, erase_regions_to_transparent

logger = logging.getLogger(__name__)


def get_vision_client_and_model():
    """從 current_app.config 取得指向 llm-proxy 的 OpenAI 相容 client 與視覺模型名。"""
    from flask import current_app
    from openai import OpenAI
    cfg = current_app.config
    base = cfg.get("OPENAI_API_BASE") or cfg.get("GOOGLE_API_BASE")
    key = cfg.get("OPENAI_API_KEY") or cfg.get("GOOGLE_API_KEY") or "sk-none"
    # 視覺用模型：優先用可走 Anthropic 視覺的 TEXT_MODEL（claude-*）
    model = cfg.get("VISION_MODEL") or cfg.get("TEXT_MODEL") or "claude-sonnet-4-6"
    client = OpenAI(api_key=key, base_url=base, timeout=180)
    return client, model


def _delete_all_slides(prs: Presentation) -> None:
    """移除母版檔內既有的範例投影片（只保留版面配置）。

    同時卸除投影片關聯，避免舊 slideN.xml part 仍被序列化而產生
    'Duplicate name: ppt/slides/slideN.xml' 警告與檔案損毀風險。
    """
    from pptx.oxml.ns import qn
    sldIdLst = prs.slides._sldIdLst
    for sldId in list(sldIdLst):
        rId = sldId.get(qn('r:id'))
        if rId:
            try:
                prs.part.drop_rel(rId)
            except Exception:
                pass
        sldIdLst.remove(sldId)


def _pick_layout(prs: Presentation):
    """挑一個乾淨、能保留母版背景的版面配置（優先『空白頁面』）。"""
    prefer = ("空白頁面", "空白", "blank", "標題及內容", "只有標題", "title only")
    layouts = list(prs.slide_layouts)
    for name in prefer:
        for lay in layouts:
            if name.lower() in (lay.name or "").lower():
                return lay
    # 退而求其次：佔位框最少者
    return min(layouts, key=lambda l: len(list(l.placeholders))) if layouts else prs.slide_layouts[0]


def _estimate_font_size(height_in: float, rtype: str) -> int:
    pt = height_in * 72 * 0.80
    if rtype == "title":
        return max(20, min(54, round(pt)))
    return max(10, min(28, round(pt)))


def _hex_to_rgb(color: Optional[str]) -> RGBColor:
    try:
        c = (color or "#333333").lstrip("#")
        return RGBColor(int(c[0:2], 16), int(c[2:4], 16), int(c[4:6], 16))
    except Exception:
        return RGBColor(0x33, 0x33, 0x33)


def _add_text_regions(slide, text_regions, title, points, sw_in, sh_in):
    """把文字放成可編輯文字框。優先用 OCR bbox；無則用 title/points 降級排版。"""
    def add_box(text, x0, y0, x1, y1, rtype, color="#333333"):
        text = (text or "").strip()
        if not text:
            return
        left = Inches(max(0.0, x0) * sw_in)
        top = Inches(max(0.0, y0) * sh_in)
        width = Inches(max(0.05, (x1 - x0)) * sw_in)
        height = Inches(max(0.03, (y1 - y0)) * sh_in)
        tb = slide.shapes.add_textbox(left, top, width, height)
        tf = tb.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.TOP
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = text
        run.font.size = Pt(_estimate_font_size((y1 - y0) * sh_in, rtype))
        run.font.bold = (rtype == "title")
        run.font.color.rgb = _hex_to_rgb(color)

    if text_regions:
        for r in text_regions:
            rtype = r.get("type", "other")
            add_box(r.get("text", ""), r.get("x0", 0), r.get("y0", 0),
                    r.get("x1", 1), r.get("y1", 0.1), rtype, r.get("color", "#333333"))
        return
    # 降級：沒有 OCR bbox 時，標題置頂、要點條列於左側
    if title:
        add_box(title, 0.07, 0.05, 0.93, 0.17, "title", "#1A1A1A")
    if points:
        body = "\n".join(f"• {p}" for p in points)
        add_box(body, 0.07, 0.22, 0.6, 0.9, "bullet", "#333333")


def _place_full_bleed(slide, pil_image, sw_in, sh_in):
    """把（可能含透明區的）底圖鋪滿整頁，置於最底層。"""
    bio = io.BytesIO()
    pil_image.save(bio, format="PNG")
    bio.seek(0)
    pic = slide.shapes.add_picture(bio, 0, 0, Inches(sw_in), Inches(sh_in))
    # 移到最底層（在文字之前加入即為底層，這裡確保順序）
    try:
        spTree = slide.shapes._spTree
        spTree.remove(pic._element)
        spTree.insert(2, pic._element)  # 2 = 前兩個為 nvGrpSpPr/grpSpPr 之後
    except Exception:
        pass
    return pic


def build_pptx_from_master(master_path: str, pages_data: List[Dict], output_file: str,
                           remove_logo: bool = True) -> Optional[str]:
    """用母版產出套版後的 PPTX（v2：保留 AI 底圖、去文字、擦除與母版重複的 logo）。

    每頁結構：
      底層 = AI「去字底圖」（bg_image_path，已擦掉文字）；其上與母版重複的
             logo/浮水印會被偵測並擦成透明，讓母版自身的 logo 透出、避免重複。
      上層 = 依 OCR bbox 放的可編輯文字框。
      母版 = 提供尺寸 / 佈景 / 頁尾 / 品牌識別。

    Args:
        master_path: 母版 .pptx / .potx 路徑
        pages_data:  [{'image_path','bg_image_path','title','points','text_regions'}...]
        output_file: 輸出路徑
        remove_logo: 是否偵測並擦除與母版重複的 logo
    Returns:
        output_file
    """
    prs = Presentation(master_path)
    sw_in = prs.slide_width / 914400.0
    sh_in = prs.slide_height / 914400.0
    _delete_all_slides(prs)
    layout = _pick_layout(prs)
    logger.info("套用母版匯出(v2)：版面=%r 尺寸=%.2fx%.2f in 共 %d 頁",
                layout.name, sw_in, sh_in, len(pages_data))

    client = model = None
    if remove_logo:
        try:
            client, model = get_vision_client_and_model()
        except Exception as e:
            logger.warning("無法取得視覺 client，將略過 logo 擦除：%s", e)
            client = None

    for idx, page in enumerate(pages_data, 1):
        slide = prs.slides.add_slide(layout)

        # 底圖：優先用去字圖（已移除文字）；無則退回原圖
        base_path = page.get("bg_image_path") or page.get("image_path")
        if base_path and os.path.exists(base_path):
            try:
                base_img = Image.open(base_path).convert("RGB")
                # 偵測並擦除與母版重複的 logo/浮水印 → 透明，讓母版 logo 透出
                if client:
                    try:
                        logos = detect_logos(base_img, client, model)
                        logger.info("第 %d 頁：偵測到 %d 個 logo/浮水印，擦除", idx, len(logos))
                        base_img = erase_regions_to_transparent(base_img, logos)
                    except Exception as e:
                        logger.warning("第 %d 頁 logo 偵測/擦除失敗（保留原底圖）：%s", idx, e)
                        base_img = base_img.convert("RGBA")
                else:
                    base_img = base_img.convert("RGBA")
                _place_full_bleed(slide, base_img, sw_in, sh_in)
            except Exception as e:
                logger.warning("第 %d 頁底圖處理失敗（略過底圖）：%s", idx, e)
        else:
            logger.warning("第 %d 頁無可用底圖：%s", idx, base_path)

        # 可編輯文字（疊在底圖之上）
        _add_text_regions(slide, page.get("text_regions"),
                          page.get("title", ""), page.get("points", []), sw_in, sh_in)

    prs.save(output_file)
    return output_file
