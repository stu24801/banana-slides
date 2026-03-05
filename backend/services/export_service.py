"""
Export Service - handles PPTX and PDF export
Based on demo.py create_pptx_from_images()
"""
import os
import json
import logging
import tempfile
from datetime import datetime, timezone
from pathlib import Path
from typing import List, Dict, Any, Optional, Tuple
from textwrap import dedent
from dataclasses import dataclass, field
from pptx import Presentation
from pptx.util import Inches
from PIL import Image
import io
import tempfile
import img2pdf
logger = logging.getLogger(__name__)


@dataclass
class ExportWarnings:
    """
    匯出過程中收集的警告資訊
    
    用於追蹤哪些操作沒有按預期執行，並反饋給前端
    """
    # 樣式提取失敗的元素
    style_extraction_failed: List[Dict[str, Any]] = field(default_factory=list)
    
    # 文字渲染失敗的元素
    text_render_failed: List[Dict[str, Any]] = field(default_factory=list)
    
    # 圖片新增失敗
    image_add_failed: List[Dict[str, Any]] = field(default_factory=list)
    
    # JSON 解析失敗（重試後仍失敗）
    json_parse_failed: List[Dict[str, Any]] = field(default_factory=list)
    
    # 其他警告
    other_warnings: List[str] = field(default_factory=list)
    
    def add_style_extraction_failed(self, element_id: str, reason: str):
        """記錄樣式提取失敗"""
        self.style_extraction_failed.append({
            'element_id': element_id,
            'reason': reason
        })
    
    def add_text_render_failed(self, text: str, reason: str):
        """記錄文字渲染失敗"""
        self.text_render_failed.append({
            'text': text[:50] + '...' if len(text) > 50 else text,
            'reason': reason
        })
    
    def add_image_failed(self, path: str, reason: str):
        """記錄圖片新增失敗"""
        self.image_add_failed.append({
            'path': path,
            'reason': reason
        })
    
    def add_json_parse_failed(self, context: str, reason: str):
        """記錄 JSON 解析失敗"""
        self.json_parse_failed.append({
            'context': context,
            'reason': reason
        })
    
    def add_warning(self, message: str):
        """新增其他警告"""
        self.other_warnings.append(message)
    
    def has_warnings(self) -> bool:
        """是否有警告"""
        return bool(
            self.style_extraction_failed or 
            self.text_render_failed or 
            self.image_add_failed or
            self.json_parse_failed or
            self.other_warnings
        )
    
    def to_summary(self) -> List[str]:
        """生成警告摘要（適合前端展示）"""
        summary = []
        
        if self.style_extraction_failed:
            summary.append(f"⚠️ {len(self.style_extraction_failed)} 個文字元素樣式提取失敗")
        
        if self.text_render_failed:
            summary.append(f"⚠️ {len(self.text_render_failed)} 個文字元素渲染失敗")
        
        if self.image_add_failed:
            summary.append(f"⚠️ {len(self.image_add_failed)} 張圖片新增失敗")
        
        if self.json_parse_failed:
            summary.append(f"⚠️ {len(self.json_parse_failed)} 次 AI 響應解析失敗")
        
        for warning in self.other_warnings[:5]:  # 最多顯示5條其他警告
            summary.append(f"⚠️ {warning}")
        
        if len(self.other_warnings) > 5:
            summary.append(f"  ...還有 {len(self.other_warnings) - 5} 條其他警告")
        
        return summary
    
    def to_dict(self) -> Dict[str, Any]:
        """轉換為字典（詳細資訊）"""
        return {
            'style_extraction_failed': self.style_extraction_failed,
            'text_render_failed': self.text_render_failed,
            'image_add_failed': self.image_add_failed,
            'json_parse_failed': self.json_parse_failed,
            'other_warnings': self.other_warnings,
            'total_warnings': (
                len(self.style_extraction_failed) + 
                len(self.text_render_failed) + 
                len(self.image_add_failed) +
                len(self.json_parse_failed) +
                len(self.other_warnings)
            )
        }


class ExportService:
    """Service for exporting presentations"""
    
    # NOTE: clean background生成功能已遷移到解耦的InpaintProvider實現
    # - DefaultInpaintProvider: 基於mask的精確區域重繪（Volcengine）
    # - GenerativeEditInpaintProvider: 基於生成式大模型的整圖編輯重繪（Gemini等）
    # 使用方式: from services.image_editability import InpaintProviderFactory
    
    @staticmethod
    def create_pptx_with_text_overlay(
        pages_data: List[Dict[str, Any]],
        output_file: str = None
    ):
        """
        建立帶有可編輯文字覆蓋層的 PPTX。
        
        每頁結構：
          底層 = 無文字背景圖（或原圖）
          上層 = 透明可編輯文字框，位置精準對齊圖片中每個文字區塊
        
        文字位置來源（優先順序）：
          1. page['text_regions'] — vision 偵測到的真實 bbox（0~1 比例座標）
          2. 降級：page['title'] 放頂部，page['points'] 放中下區域
        
        Args:
            pages_data: 每頁資料：
                [{
                  "image_path": str,
                  "title": str,
                  "points": [str, ...],
                  "text_regions": [  ← 可選，vision 偵測結果
                    {"text": str, "x0": float, "y0": float, "x1": float, "y1": float, "type": str},
                    ...
                  ]
                }, ...]
            output_file: 輸出路徑（None → 回傳 bytes）
        """
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN

        SLIDE_W_IN = 10.0
        SLIDE_H_IN = 5.625
        SLIDE_W = Inches(SLIDE_W_IN)
        SLIDE_H = Inches(SLIDE_H_IN)

        prs = Presentation()
        try:
            core = prs.core_properties
            now = datetime.now(timezone.utc)
            core.author = "banana-slides"
            core.last_modified_by = "banana-slides"
            core.created = now
            core.modified = now
        except Exception:
            pass

        prs.slide_width = SLIDE_W
        prs.slide_height = SLIDE_H

        def _set_transparent(shape):
            shape.fill.background()
            shape.line.fill.background()

        def _add_text_box_inches(slide, left, top, width, height,
                                  text, font_size_pt, bold=False,
                                  color_rgb=(255, 255, 255),
                                  align=PP_ALIGN.LEFT):
            """新增透明文字框，座標單位為 inch"""
            tb = slide.shapes.add_textbox(
                Inches(left), Inches(top),
                Inches(width), Inches(height)
            )
            _set_transparent(tb)
            tf = tb.text_frame
            tf.word_wrap = True
            tf.margin_left = Inches(0.05)
            tf.margin_right = Inches(0.05)
            tf.margin_top = Inches(0.02)
            tf.margin_bottom = Inches(0.02)
            p = tf.paragraphs[0]
            p.alignment = align
            run = p.add_run()
            run.text = text
            run.font.size = Pt(font_size_pt)
            run.font.bold = bold
            run.font.color.rgb = RGBColor(*color_rgb)
            return tb

        def _region_to_inches(region):
            """將 0~1 比例 bbox 轉為 inch"""
            x0 = region['x0'] * SLIDE_W_IN
            y0 = region['y0'] * SLIDE_H_IN
            x1 = region['x1'] * SLIDE_W_IN
            y1 = region['y1'] * SLIDE_H_IN
            return x0, y0, max(0.3, x1 - x0), max(0.2, y1 - y0)

        def _estimate_font_size(height_in, rtype):
            """根據區塊高度估算字型大小"""
            pt = height_in * 72 * 0.65  # 72pt/inch，65% 填充率
            if rtype == 'title':
                return max(18, min(48, round(pt)))
            else:
                return max(10, min(28, round(pt)))

        for page in pages_data:
            image_path = page.get('image_path', '')
            title = page.get('title', '')
            points = page.get('points', [])
            text_regions = page.get('text_regions')  # vision 偵測結果

            blank_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(blank_layout)

            # ── 底層：背景圖 ───────────────────────────────────────────────────
            if image_path and os.path.exists(image_path):
                slide.shapes.add_picture(image_path, 0, 0, SLIDE_W, SLIDE_H)
            else:
                logger.warning(f"圖片不存在，跳過：{image_path}")

            # ── 上層：精準文字框 ───────────────────────────────────────────────
            if text_regions:
                # ✅ 模式 A：使用 vision 偵測的真實 bbox
                for region in text_regions:
                    rtype = region.get('type', 'other')
                    rtext = region.get('text', '').strip()
                    if not rtext:
                        continue

                    left, top, w, h = _region_to_inches(region)
                    font_size = _estimate_font_size(h, rtype)
                    bold = rtype == 'title'
                    # 字色白色（主色），label/small 用半透明白
                    color = (255, 255, 255)

                    _add_text_box_inches(
                        slide, left, top, w, h,
                        text=rtext,
                        font_size_pt=font_size,
                        bold=bold,
                        color_rgb=color,
                        align=PP_ALIGN.LEFT
                    )
                    logger.debug(f"Added region [{rtype}] '{rtext[:30]}' at ({left:.2f},{top:.2f})")

            else:
                # ⚠️ 模式 B：降級 — 用 outline_content 的 title / points，位置固定
                logger.debug("No text_regions, using fallback title/points layout")
                if title:
                    _add_text_box_inches(
                        slide, 0.3, 0.2, 9.4, 1.1,
                        text=title, font_size_pt=32, bold=True,
                        color_rgb=(255, 255, 255), align=PP_ALIGN.LEFT
                    )
                if points:
                    bullets_text = '\n'.join(
                        f'• {pt}' if not pt.startswith('•') else pt
                        for pt in points
                    )
                    _add_text_box_inches(
                        slide, 0.4, 1.5, 9.2, 3.8,
                        text=bullets_text, font_size_pt=18, bold=False,
                        color_rgb=(255, 255, 255), align=PP_ALIGN.LEFT
                    )

        if output_file:
            prs.save(output_file)
            return None
        else:
            buf = io.BytesIO()
            prs.save(buf)
            buf.seek(0)
            return buf.getvalue()

    @staticmethod
    def create_pptx_from_images(image_paths: List[str], output_file: str = None) -> bytes:
        """
        Create PPTX file from image paths
        Based on demo.py create_pptx_from_images()
        
        Args:
            image_paths: List of absolute paths to images
            output_file: Optional output file path (if None, returns bytes)
        
        Returns:
            PPTX file as bytes if output_file is None
        """
        # Create presentation
        prs = Presentation()
        
        # Set author/date metadata for exported PPTX
        try:
            core = prs.core_properties
            now = datetime.now(timezone.utc)
            core.author = "banana-slides"
            core.last_modified_by = "banana-slides"
            core.created = now
            core.modified = now
            core.last_printed = None
        except Exception as e:
            logger.warning(f"Failed to set core properties: {e}")
        
        # Set slide dimensions to 16:9 (width 10 inches, height 5.625 inches)
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(5.625)
        
        # Add each image as a slide
        for image_path in image_paths:
            if not os.path.exists(image_path):
                logger.warning(f"Image not found: {image_path}")
                continue
            
            # Add blank slide layout (layout 6 is typically blank)
            blank_slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(blank_slide_layout)
            
            # Add image to fill entire slide
            slide.shapes.add_picture(
                image_path,
                left=0,
                top=0,
                width=prs.slide_width,
                height=prs.slide_height
            )
        
        # Save or return bytes
        if output_file:
            prs.save(output_file)
            return None
        else:
            # Save to bytes
            pptx_bytes = io.BytesIO()
            prs.save(pptx_bytes)
            pptx_bytes.seek(0)
            return pptx_bytes.getvalue()
    
    @staticmethod
    def create_pdf_from_images(image_paths: List[str], output_file: str = None) -> Optional[bytes]:
        """
        Create PDF file from image paths using img2pdf (low memory usage)

        Args:
            image_paths: List of absolute paths to images
            output_file: Optional output file path (if None, returns bytes)

        Returns:
            PDF file as bytes if output_file is None, otherwise None
        """
        # Validate images exist and log warnings for missing files
        valid_paths = []
        for p in image_paths:
            if os.path.exists(p):
                valid_paths.append(p)
            else:
                logger.warning(f"Image not found and will be skipped for PDF export: {p}")

        if not valid_paths:
            raise ValueError("No valid images found for PDF export")

        try:
            logger.info(f"Using img2pdf for PDF export ({len(valid_paths)} pages, low memory mode)")

            # Set page layout: 16:9 aspect ratio (10 inches × 5.625 inches)
            layout_fun = img2pdf.get_layout_fun(
                pagesize=(img2pdf.in_to_pt(10), img2pdf.in_to_pt(5.625))
            )

            # Convert images to PDF
            pdf_bytes = img2pdf.convert(valid_paths, layout_fun=layout_fun)

            if output_file:
                with open(output_file, "wb") as f:
                    f.write(pdf_bytes)
                return None
            else:
                return pdf_bytes
        except (img2pdf.ImageOpenError, ValueError, IOError) as e:
            logger.warning(f"img2pdf conversion failed: {e}. Falling back to Pillow (high memory usage).")
            return ExportService.create_pdf_from_images_pillow(valid_paths, output_file)

    @staticmethod
    def create_pdf_from_images_pillow(image_paths: List[str], output_file: str = None) -> Optional[bytes]:
        """
        Create PDF file from image paths using Pillow (original method)

        Note: This method loads all images into memory at once.
        For large projects (50+ pages with 20MB/page), use create_pdf_from_images instead.

        Args:
            image_paths: List of absolute paths to images
            output_file: Optional output file path (if None, returns bytes)

        Returns:
            PDF file as bytes if output_file is None, otherwise None
        """
        images = []

        # Load all images
        for image_path in image_paths:
            if not os.path.exists(image_path):
                logger.warning(f"Image not found: {image_path}")
                continue

            img = Image.open(image_path)

            # Convert to RGB if necessary (PDF requires RGB)
            if img.mode != 'RGB':
                img = img.convert('RGB')

            images.append(img)

        if not images:
            raise ValueError("No valid images found for PDF export")

        # Save as PDF
        if output_file:
            images[0].save(
                output_file,
                save_all=True,
                append_images=images[1:],
                format='PDF'
            )
            return None
        else:
            # Save to bytes
            pdf_bytes = io.BytesIO()
            images[0].save(
                pdf_bytes,
                save_all=True,
                append_images=images[1:],
                format='PDF'
            )
            pdf_bytes.seek(0)
            return pdf_bytes.getvalue()
       
    @staticmethod
    def _add_mineru_text_to_slide(builder, slide, text_item: Dict[str, Any], scale_x: float = 1.0, scale_y: float = 1.0):
        """
        Add text item from MinerU to slide
        
        Args:
            builder: PPTXBuilder instance
            slide: Target slide
            text_item: Text item from MinerU content_list
            scale_x: X-axis scale factor
            scale_y: Y-axis scale factor
        """
        text = text_item.get('text', '').strip()
        if not text:
            return
        
        bbox = text_item.get('bbox')
        if not bbox or len(bbox) != 4:
            logger.warning(f"Invalid bbox for text item: {text_item}")
            return
        
        original_bbox = bbox.copy()
        
        # Apply scale factors to bbox
        x0, y0, x1, y1 = bbox
        bbox = [
            int(x0 * scale_x),
            int(y0 * scale_y),
            int(x1 * scale_x),
            int(y1 * scale_y)
        ]
        
        if scale_x != 1.0 or scale_y != 1.0:
            logger.debug(f"Text bbox scaled: {original_bbox} -> {bbox} (scale: {scale_x:.3f}x{scale_y:.3f})")
        
        # Determine text level (only used for styling like bold, NOT for font size)
        # Font size is purely calculated from bbox dimensions
        item_type = text_item.get('type', 'text')
        text_level = text_item.get('text_level')
        
        # Map to level for styling purposes (bold titles)
        if item_type == 'title' or text_level == 1:
            level = 'title'  # Will be bold
        else:
            level = 'default'
        
        # Add text element
        # Note: text_level is only used for bold styling, not font size calculation
        try:
            builder.add_text_element(
                slide=slide,
                text=text,
                bbox=bbox,
                text_level=level  # For styling (bold) only, not font size
            )
        except Exception as e:
            logger.error(f"Failed to add text element: {str(e)}")
    
    @staticmethod
    def _add_table_cell_elements_to_slide(
        builder,
        slide,
        cell_elements: List[Dict[str, Any]],
        scale_x: float = 1.0,
        scale_y: float = 1.0
    ):
        """
        Add table cell elements as individual text boxes to slide
        這些單元格元素已經有正確的全域性bbox座標
        
        Args:
            builder: PPTXBuilder instance
            slide: Target slide
            cell_elements: List of EditableElement (table_cell type)
            scale_x: X-axis scale factor
            scale_y: Y-axis scale factor
        """
        from pptx.util import Pt
        from pptx.dml.color import RGBColor
        
        logger.info(f"開始新增表格單元格元素，共 {len(cell_elements)} 個")
        
        for cell_elem in cell_elements:
            text = cell_elem.get('content', '')
            bbox_global = cell_elem.get('bbox_global', {})
            
            if not text.strip():
                continue
            
            # bbox_global已經是全域性座標，直接使用並應用縮放
            x0 = bbox_global.get('x0', 0)
            y0 = bbox_global.get('y0', 0)
            x1 = bbox_global.get('x1', 0)
            y1 = bbox_global.get('y1', 0)
            
            # 構建bbox列表 [x0, y0, x1, y1] 並應用縮放
            bbox = [
                int(x0 * scale_x),
                int(y0 * scale_y),
                int(x1 * scale_x),
                int(y1 * scale_y)
            ]
            
            try:
                # 使用已有的 add_text_element 方法新增文字框（不新增邊框）
                builder.add_text_element(
                    slide=slide,
                    text=text,
                    bbox=bbox,
                    text_level=None,
                    align='center'
                )
                
                logger.debug(f"  新增單元格: '{text[:10]}...' at bbox {bbox}")
                
            except Exception as e:
                logger.warning(f"新增單元格失敗: {e}")
        
        logger.info(f"✓ 表格單元格新增完成，共 {len(cell_elements)} 個")
    
    @staticmethod
    def _add_mineru_image_to_slide(
        builder,
        slide,
        image_item: Dict[str, Any],
        mineru_dir: Path,
        scale_x: float = 1.0,
        scale_y: float = 1.0
    ):
        """
        Add image or table item from MinerU to slide
        
        Args:
            builder: PPTXBuilder instance
            slide: Target slide
            image_item: Image/table item from MinerU content_list
            mineru_dir: MinerU result directory
            scale_x: X-axis scale factor
            scale_y: Y-axis scale factor
        """
        bbox = image_item.get('bbox')
        if not bbox or len(bbox) != 4:
            logger.warning(f"Invalid bbox for image item: {image_item}")
            return
        
        original_bbox = bbox.copy()
        
        # Apply scale factors to bbox
        x0, y0, x1, y1 = bbox
        bbox = [
            int(x0 * scale_x),
            int(y0 * scale_y),
            int(x1 * scale_x),
            int(y1 * scale_y)
        ]
        
        if scale_x != 1.0 or scale_y != 1.0:
            logger.debug(f"Item bbox scaled: {original_bbox} -> {bbox} (scale: {scale_x:.3f}x{scale_y:.3f})")
        
        # Check if this is a table with子元素 (cells from Baidu OCR)
        item_type = image_item.get('element_type') or image_item.get('type', 'image')
        children = image_item.get('children', [])
        
        logger.debug(f"Processing {item_type} element, has {len(children)} children")
        
        if children and item_type == 'table':
            # Add editable table from child elements (cells)
            try:
                # Filter only table_cell elements
                cell_elements = [child for child in children if child.get('element_type') == 'table_cell']
                
                if cell_elements:
                    logger.info(f"新增可編輯表格（{len(cell_elements)}個單元格）")
                    ExportService._add_table_cell_elements_to_slide(
                        builder=builder,
                        slide=slide,
                        cell_elements=cell_elements,
                        scale_x=scale_x,
                        scale_y=scale_y
                    )
                    return  # Table added successfully
            except Exception as e:
                logger.exception("Failed to add table cells, falling back to image")
                # Fall through to add as image instead
        
        # Check if this is a table with HTML data (legacy)
        html_table = image_item.get('html_table')
        if html_table and item_type == 'table':
            # Add editable table from HTML
            try:
                builder.add_table_element(
                    slide=slide,
                    html_table=html_table,
                    bbox=bbox
                )
                logger.info(f"Added editable table at bbox {bbox}")
                return  # Table added successfully
            except Exception as e:
                logger.error(f"Failed to add table: {str(e)}, falling back to image")
                # Fall through to add as image instead
        
        # Add as image (either image type or table fallback)
        img_path_str = image_item.get('img_path', '')
        if not img_path_str:
            logger.warning(f"No img_path in item: {image_item}")
            return
        
        # Try to find the image file
        # MinerU may store images in 'images/' subdirectory
        possible_paths = [
            mineru_dir / img_path_str,
            mineru_dir / 'images' / Path(img_path_str).name,
            mineru_dir / Path(img_path_str).name,
        ]
        
        image_path = None
        for path in possible_paths:
            if path.exists():
                image_path = str(path)
                break
        
        if not image_path:
            logger.warning(f"Image file not found: {img_path_str}")
            # Add placeholder
            builder.add_image_placeholder(slide, bbox)
            return
        
        # Add image element
        try:
            builder.add_image_element(
                slide=slide,
                image_path=image_path,
                bbox=bbox
            )
        except Exception as e:
            logger.error(f"Failed to add image element: {str(e)}")
    
    @staticmethod
    def _collect_text_elements_for_extraction(
        elements: List,  # List[EditableElement]
        depth: int = 0
    ) -> List[tuple]:
        """
        遞迴收集所有需要提取樣式的文字元素
        
        Args:
            elements: EditableElement列表
            depth: 當前遞迴深度
        
        Returns:
            元組列表，每個元組為 (element_id, image_path, text_content)
        """
        text_items = []
        
        for elem in elements:
            elem_type = elem.element_type
            
            # 文字型別元素需要提取樣式
            if elem_type in ['text', 'title', 'table_cell', 'list', 'paragraph', 'header', 'footer', 'heading', 'table_caption', 'image_caption']:
                if elem.content and elem.image_path and os.path.exists(elem.image_path):
                    text = elem.content.strip()
                    if text:
                        text_items.append((elem.element_id, elem.image_path, text))
            
            # 遞迴處理子元素
            if hasattr(elem, 'children') and elem.children:
                child_items = ExportService._collect_text_elements_for_extraction(
                    elements=elem.children,
                    depth=depth + 1
                )
                text_items.extend(child_items)
        
        return text_items
    
    @staticmethod
    def _batch_extract_text_styles(
        text_items: List[tuple],
        text_attribute_extractor,
        max_workers: int = 8
    ) -> Dict[str, Any]:
        """
        批次並行提取文字樣式（逐個裁剪區域分析）
        
        此方法對每一段文字的裁剪區域單獨進行分析。
        經測試，此方法效果較好，目前仍在使用。
        
        備選方案：_batch_extract_text_styles_with_full_image 可一次性分析全圖所有文字。
        
        Args:
            text_items: 元組列表，每個元組為 (element_id, image_path, text_content)
            text_attribute_extractor: 文字屬性提取器
            max_workers: 併發數
        
        Returns:
            字典，key為element_id，value為TextStyleResult
        """
        from concurrent.futures import ThreadPoolExecutor, as_completed
        
        if not text_items or not text_attribute_extractor:
            return {}
        
        logger.info(f"並行提取 {len(text_items)} 個文字元素的樣式（併發數: {max_workers}）...")
        
        results = {}
        
        def extract_single(item):
            element_id, image_path, text_content = item
            try:
                style = text_attribute_extractor.extract(
                    image=image_path,
                    text_content=text_content
                )
                return element_id, style
            except Exception as e:
                logger.warning(f"提取文字樣式失敗 [{element_id}]: {e}")
                return element_id, None
        
        with ThreadPoolExecutor(max_workers=max_workers) as executor:
            futures = {executor.submit(extract_single, item): item[0] for item in text_items}
            
            for future in as_completed(futures):
                element_id, style = future.result()
                if style is not None:
                    results[element_id] = style
        
        logger.info(f"✓ 文字樣式提取完成，成功 {len(results)}/{len(text_items)} 個")
        return results
    
    @staticmethod
    def _collect_text_elements_for_batch_extraction(
        elements: List,  # List[EditableElement]
        depth: int = 0
    ) -> List[Dict[str, Any]]:
        """
        遞迴收集所有需要批次提取樣式的文字元素（新格式，包含bbox）
        
        Args:
            elements: EditableElement列表
            depth: 當前遞迴深度
        
        Returns:
            字典列表，每個字典包含 element_id, bbox, content
        """
        text_items = []
        
        for elem in elements:
            elem_type = elem.element_type
            
            # 文字型別元素需要提取樣式
            if elem_type in ['text', 'title', 'table_cell', 'list', 'paragraph', 'header', 'footer', 'heading', 'table_caption', 'image_caption']:
                if elem.content:
                    text = elem.content.strip()
                    if text:
                        # 使用全域性座標 bbox_global
                        bbox = elem.bbox_global if hasattr(elem, 'bbox_global') and elem.bbox_global else elem.bbox
                        text_items.append({
                            'element_id': elem.element_id,
                            'bbox': [bbox.x0, bbox.y0, bbox.x1, bbox.y1],
                            'content': text
                        })
            
            # 遞迴處理子元素
            if hasattr(elem, 'children') and elem.children:
                child_items = ExportService._collect_text_elements_for_batch_extraction(
                    elements=elem.children,
                    depth=depth + 1
                )
                text_items.extend(child_items)
        
        return text_items
    
    @staticmethod
    def _batch_extract_text_styles_with_full_image(
        editable_images: List,  # List[EditableImage]
        text_attribute_extractor,
        max_workers: int = 4
    ) -> Dict[str, Any]:
        """
        【新邏輯】使用全圖批次提取所有文字樣式
        
        新方法：給 caption model 提供全圖，以及提取後的所有文字 bbox 和內容，
        讓模型一次性分析所有文字的樣式屬性（顏色、粗體、對齊等）。
        
        優勢：模型可以看到全域性資訊，分析更準確。
        
        Args:
            editable_images: EditableImage列表，每個對應一張PPT頁面
            text_attribute_extractor: 文字屬性提取器（需要有 extract_batch_with_full_image 方法）
            max_workers: 併發處理頁面數
        
        Returns:
            字典，key為element_id，value為TextStyleResult
        """
        from concurrent.futures import ThreadPoolExecutor, as_completed
        
        if not editable_images or not text_attribute_extractor:
            return {}
        
        # 檢查提取器是否支援批次提取
        if not hasattr(text_attribute_extractor, 'extract_batch_with_full_image'):
            logger.warning("提取器不支援 extract_batch_with_full_image 方法，回退到舊邏輯")
            # 回退到舊邏輯
            all_text_items = []
            for editable_img in editable_images:
                text_items = ExportService._collect_text_elements_for_extraction(editable_img.elements)
                all_text_items.extend(text_items)
            return ExportService._batch_extract_text_styles(
                text_items=all_text_items,
                text_attribute_extractor=text_attribute_extractor,
                max_workers=max_workers * 2
            )
        
        logger.info(f"【新邏輯】使用全圖批次分析 {len(editable_images)} 頁的文字樣式...")
        
        all_results = {}
        
        def process_single_page(editable_img, page_idx):
            """處理單個頁面的文字樣式提取"""
            try:
                # 收集該頁面的所有文字元素
                text_elements = ExportService._collect_text_elements_for_batch_extraction(
                    editable_img.elements
                )
                
                if not text_elements:
                    logger.info(f"  頁面 {page_idx + 1}: 無文字元素")
                    return {}
                
                logger.info(f"  頁面 {page_idx + 1}: 分析 {len(text_elements)} 個文字元素...")
                
                # 使用原始圖片路徑作為全圖
                full_image_path = editable_img.image_path
                
                # 呼叫批次提取方法
                page_results = text_attribute_extractor.extract_batch_with_full_image(
                    full_image=full_image_path,
                    text_elements=text_elements
                )
                
                logger.info(f"  頁面 {page_idx + 1}: 成功提取 {len(page_results)} 個元素的樣式")
                return page_results
                
            except Exception as e:
                logger.error(f"頁面 {page_idx + 1} 文字樣式提取失敗: {e}", exc_info=True)
                return {}
        
        # 併發處理所有頁面
        with ThreadPoolExecutor(max_workers=max_workers) as executor:
            futures = {
                executor.submit(process_single_page, img, idx): idx 
                for idx, img in enumerate(editable_images)
            }
            
            for future in as_completed(futures):
                page_idx = futures[future]
                try:
                    page_results = future.result()
                    all_results.update(page_results)
                except Exception as e:
                    logger.error(f"頁面 {page_idx + 1} 處理失敗: {e}")
        
        total_elements = sum(
            len(ExportService._collect_text_elements_for_batch_extraction(img.elements))
            for img in editable_images
        )
        logger.info(f"✓ 全圖批次文字樣式提取完成，成功 {len(all_results)}/{total_elements} 個")
        
        return all_results
    
    @staticmethod
    def _batch_extract_text_styles_hybrid(
        editable_images: List,  # List[EditableImage]
        text_attribute_extractor,
        max_workers: int = 8
    ) -> Tuple[Dict[str, Any], List[Tuple[str, str]]]:
        """
        【混合策略】結合全域性識別和單個裁剪識別的優勢
        
        策略：
        - 全域性識別（全圖分析）：獲取 is_bold、is_italic、is_underline、text_alignment
          因為這些屬性需要看整體佈局和上下文才能判斷準確
        - 單個裁剪識別：獲取 font_color
          因為顏色需要精確看區域性畫素才能識別準確
        
        Args:
            editable_images: EditableImage列表，每個對應一張PPT頁面
            text_attribute_extractor: 文字屬性提取器
            max_workers: 併發數
        
        Returns:
            (results, failed_extractions):
            - results: 字典，key為element_id，value為TextStyleResult（合併後的結果）
            - failed_extractions: 失敗列表，每項為 (element_id, error_reason)
        """
        from concurrent.futures import ThreadPoolExecutor, as_completed
        from services.image_editability.text_attribute_extractors import TextStyleResult
        
        if not editable_images or not text_attribute_extractor:
            return {}, []
        
        # 檢查提取器是否支援批次提取
        if not hasattr(text_attribute_extractor, 'extract_batch_with_full_image'):
            logger.warning("提取器不支援混合策略，回退到單個裁剪識別")
            all_text_items = []
            for editable_img in editable_images:
                text_items = ExportService._collect_text_elements_for_extraction(editable_img.elements)
                all_text_items.extend(text_items)
            results = ExportService._batch_extract_text_styles(
                text_items=all_text_items,
                text_attribute_extractor=text_attribute_extractor,
                max_workers=max_workers
            )
            return results, []  # 回退方法暫不收集失敗資訊
        
        logger.info(f"【混合策略】開始分析 {len(editable_images)} 頁的文字樣式...")
        logger.info(f"  - 全域性識別: is_bold, is_italic, is_underline, text_alignment")
        logger.info(f"  - 單個識別: font_color")
        
        # Step 1: 收集所有文字元素
        all_text_items = []  # 用於單個裁剪識別 (element_id, image_path, content)
        page_text_elements = {}  # 用於全域性識別 {page_idx: [text_elements]}
        
        for page_idx, editable_img in enumerate(editable_images):
            # 收集用於單個裁剪識別的資料
            text_items = ExportService._collect_text_elements_for_extraction(editable_img.elements)
            all_text_items.extend(text_items)
            
            # 收集用於全域性識別的資料
            batch_elements = ExportService._collect_text_elements_for_batch_extraction(editable_img.elements)
            if batch_elements:
                page_text_elements[page_idx] = {
                    'image_path': editable_img.image_path,
                    'elements': batch_elements
                }
        
        if not all_text_items:
            return {}
        
        # Step 2: 並行執行兩種識別
        global_results = {}  # 全域性識別結果
        local_results = {}   # 單個裁剪識別結果
        
        def extract_global_for_page(page_idx, page_data):
            """全域性識別單頁"""
            try:
                results = text_attribute_extractor.extract_batch_with_full_image(
                    full_image=page_data['image_path'],
                    text_elements=page_data['elements']
                )
                return page_idx, results
            except Exception as e:
                logger.warning(f"全域性識別頁面 {page_idx + 1} 失敗: {e}")
                return page_idx, {}
        
        # 收集失敗資訊
        failed_extractions = []  # [(element_id, reason), ...]
        
        def extract_local_single(item):
            """單個裁剪識別"""
            element_id, image_path, text_content = item
            try:
                style = text_attribute_extractor.extract(
                    image=image_path,
                    text_content=text_content
                )
                # 只要 style 不為 None 就算成功（黑色也是有效顏色）
                if style:
                    return element_id, style, None
                else:
                    return element_id, None, "樣式提取返回空"
            except Exception as e:
                logger.warning(f"單個識別失敗 [{element_id}]: {e}")
                return element_id, None, str(e)
        
        # 併發執行全域性識別和單個裁剪識別
        logger.info(f"  併發執行: 全域性識別 {len(page_text_elements)} 頁 + 單個識別 {len(all_text_items)} 個元素...")
        
        with ThreadPoolExecutor(max_workers=max_workers) as executor:
            # 提交全域性識別任務
            global_futures = {
                executor.submit(extract_global_for_page, idx, data): ('global', idx)
                for idx, data in page_text_elements.items()
            }
            
            # 提交單個裁剪識別任務
            local_futures = {
                executor.submit(extract_local_single, item): ('local', item[0])
                for item in all_text_items
            }
            
            # 收集全域性識別結果
            for future in as_completed(global_futures):
                task_type, page_idx = global_futures[future]
                try:
                    _, page_results = future.result()
                    global_results.update(page_results)
                except Exception as e:
                    logger.error(f"全域性識別任務失敗: {e}")
            
            # 收集單個裁剪識別結果
            for future in as_completed(local_futures):
                task_type, element_id = local_futures[future]
                try:
                    elem_id, style, error = future.result()
                    if style is not None:
                        local_results[elem_id] = style
                    if error:
                        failed_extractions.append((elem_id, error))
                except Exception as e:
                    logger.error(f"單個識別任務失敗: {e}")
                    failed_extractions.append((element_id, str(e)))
        
        # Step 3: 合併結果
        # 優先使用全域性識別的佈局屬性，使用單個識別的顏色屬性
        merged_results = {}
        
        all_element_ids = set(global_results.keys()) | set(local_results.keys())
        
        for element_id in all_element_ids:
            global_style = global_results.get(element_id)
            local_style = local_results.get(element_id)
            
            if global_style and local_style:
                # 混合：顏色用單個識別（包括 colored_segments），佈局用全域性識別
                merged_results[element_id] = TextStyleResult(
                    font_color_rgb=local_style.font_color_rgb,  # 單個識別的顏色
                    colored_segments=local_style.colored_segments,  # 單個識別的多顏色片段
                    is_bold=global_style.is_bold,              # 全域性識別的粗體
                    is_italic=global_style.is_italic,          # 全域性識別的斜體
                    is_underline=global_style.is_underline,    # 全域性識別的下劃線
                    text_alignment=global_style.text_alignment, # 全域性識別的對齊
                    confidence=0.9,
                    metadata={
                        'source': 'hybrid',
                        'color_source': 'local',
                        'layout_source': 'global'
                    }
                )
            elif local_style:
                # 只有單個識別結果
                merged_results[element_id] = local_style
            elif global_style:
                # 只有全域性識別結果
                merged_results[element_id] = global_style
        
        logger.info(f"✓ 混合策略完成: 全域性識別 {len(global_results)} 個, 單個識別 {len(local_results)} 個, 合併 {len(merged_results)} 個, 失敗 {len(failed_extractions)} 個")
        
        return merged_results, failed_extractions
    
    @staticmethod
    def create_editable_pptx_with_recursive_analysis(
        image_paths: List[str] = None,
        output_file: str = None,
        slide_width_pixels: int = 1920,
        slide_height_pixels: int = 1080,
        max_depth: int = 2,
        max_workers: int = 8,
        editable_images: List = None,  # 可選：直接傳入已分析的EditableImage列表
        text_attribute_extractor = None,  # 可選：文字屬性提取器，用於提取顏色、粗體、斜體等樣式
        progress_callback = None,  # 可選：進度回撥函式 (step, message, percent) -> None
        export_extractor_method: str = 'hybrid',  # 元件提取方法: mineru, hybrid
        export_inpaint_method: str = 'hybrid'  # 背景修復方法: generative, baidu, hybrid
    ) -> Tuple[Optional[bytes], ExportWarnings]:
        """
        使用遞迴圖片可編輯化服務建立可編輯PPTX
        
        這是新的架構方法，使用ImageEditabilityService進行遞迴版面分析。
        
        兩種使用方式：
        1. 傳入 image_paths：自動分析圖片並生成PPTX
        2. 傳入 editable_images：直接使用已分析的結果（避免重複分析）
        
        配置（如 MinerU token）自動從 Flask app.config 獲取。
        
        Args:
            image_paths: 圖片路徑列表（可選，與editable_images二選一）
            output_file: 輸出檔案路徑（可選）
            slide_width_pixels: 目標幻燈片寬度
            slide_height_pixels: 目標幻燈片高度
            max_depth: 最大遞迴深度
            max_workers: 併發處理數
            editable_images: 已分析的EditableImage列表（可選，與image_paths二選一）
            text_attribute_extractor: 文字屬性提取器（可選），用於提取文字顏色、粗體、斜體等樣式
                可透過 TextAttributeExtractorFactory.create_caption_model_extractor() 建立
            export_extractor_method: 元件提取方法 ('mineru' 或 'hybrid'，預設 'hybrid')
            export_inpaint_method: 背景修復方法 ('generative', 'baidu', 'hybrid'，預設 'hybrid')
        
        Returns:
            (pptx_bytes, warnings): 元組，包含 PPTX 位元組流和警告資訊
            - pptx_bytes: PPTX 檔案位元組流（如果 output_file 為 None），否則為 None
            - warnings: ExportWarnings 物件，包含所有警告資訊
        """
        from services.image_editability import ServiceConfig, ImageEditabilityService
        from utils.pptx_builder import PPTXBuilder
        
        # 初始化警告收集器
        warnings = ExportWarnings()
        
        # 輔助函式：報告進度
        def report_progress(step: str, message: str, percent: int):
            logger.info(f"[進度 {percent}%] {step}: {message}")
            if progress_callback:
                try:
                    progress_callback(step, message, percent)
                except Exception as e:
                    logger.warning(f"進度回撥失敗: {e}")
        
        # 如果已提供分析結果，直接使用；否則需要分析
        if editable_images is not None:
            logger.info(f"使用已提供的 {len(editable_images)} 個分析結果建立PPTX")
            report_progress("準備", f"使用已有分析結果（{len(editable_images)} 頁）", 10)
        else:
            if not image_paths:
                raise ValueError("必須提供 image_paths 或 editable_images 之一")
            
            total_pages = len(image_paths)
            logger.info(f"開始使用遞迴分析方法建立可編輯PPTX，共 {total_pages} 頁")
            report_progress("開始", f"準備分析 {total_pages} 頁幻燈片...", 0)
            
            # 1. 建立ImageEditabilityService（配置自動從 Flask config 獲取，使用專案匯出設定）
            logger.info(f"使用匯出設定: extractor={export_extractor_method}, inpaint={export_inpaint_method}")
            config = ServiceConfig.from_defaults(
                max_depth=max_depth,
                extractor_method=export_extractor_method,
                inpaint_method=export_inpaint_method
            )
            editability_service = ImageEditabilityService(config)
            
            # 2. 併發處理所有頁面，生成EditableImage結構
            report_progress("版面分析", f"開始分析 {total_pages} 張圖片（併發數: {max_workers}）...", 5)
            from concurrent.futures import ThreadPoolExecutor, as_completed
            
            editable_images = []
            completed_count = 0
            with ThreadPoolExecutor(max_workers=max_workers) as executor:
                futures = {
                    executor.submit(editability_service.make_image_editable, img_path): idx
                    for idx, img_path in enumerate(image_paths)
                }
                
                results = [None] * len(image_paths)
                for future in as_completed(futures):
                    idx = futures[future]
                    try:
                        results[idx] = future.result()
                        completed_count += 1
                        # 版面分析佔 5% - 40% 的進度
                        percent = 5 + int(35 * completed_count / total_pages)
                        report_progress("版面分析", f"已完成第 {completed_count}/{total_pages} 頁的版面分析", percent)
                    except Exception as e:
                        logger.error(f"處理圖片 {image_paths[idx]} 失敗: {e}")
                        raise
                
                editable_images = results
        
        # 2.5. 使用混合策略提取所有文字元素的樣式（如果提供了提取器）
        # 混合策略：全域性識別（粗體/斜體/下劃線/對齊）+ 單個裁剪識別（顏色）
        text_styles_cache = {}
        if text_attribute_extractor:
            report_progress("樣式提取", "開始提取文字樣式（混合策略）...", 45)
            
            # 統計文字元素數量
            total_text_count = sum(
                len(ExportService._collect_text_elements_for_extraction(img.elements))
                for img in editable_images
            )
            
            if total_text_count > 0:
                report_progress("樣式提取", f"混合策略分析 {total_text_count} 個文字元素...", 50)
                text_styles_cache, failed_extractions = ExportService._batch_extract_text_styles_hybrid(
                    editable_images=editable_images,
                    text_attribute_extractor=text_attribute_extractor,
                    max_workers=max_workers * 2
                )
                
                # 記錄樣式提取失敗的元素（詳細）
                for element_id, reason in failed_extractions:
                    warnings.add_style_extraction_failed(element_id, reason)
                
                # 記錄彙總資訊
                extracted_count = len(text_styles_cache)
                failed_count = len(failed_extractions)
                if failed_count > 0:
                    logger.warning(f"樣式提取: {failed_count}/{total_text_count} 個元素失敗")
                
                report_progress("樣式提取", f"✓ 完成 {extracted_count}/{total_text_count} 個文字樣式提取（{failed_count} 個失敗）", 70)
        
        report_progress("構建PPTX", "開始構建可編輯PPTX檔案...", 75)
        
        # 4. 建立PPTX構建器
        builder = PPTXBuilder()
        builder.create_presentation()
        builder.setup_presentation_size(slide_width_pixels, slide_height_pixels)
        
        # 5. 為每個頁面構建幻燈片
        total_pages = len(editable_images)
        for page_idx, editable_img in enumerate(editable_images):
            # 構建PPTX佔 75% - 95% 的進度
            percent = 75 + int(20 * page_idx / total_pages)
            report_progress("構建PPTX", f"構建第 {page_idx + 1}/{total_pages} 頁...", percent)
            logger.info(f"  構建第 {page_idx + 1}/{total_pages} 頁...")
            
            # 建立空白幻燈片
            slide = builder.add_blank_slide()
            
            # 新增背景圖（參考原實現，使用slide.shapes.add_picture）
            if editable_img.clean_background and os.path.exists(editable_img.clean_background):
                logger.info(f"    新增clean background: {editable_img.clean_background}")
                try:
                    slide.shapes.add_picture(
                        editable_img.clean_background,
                        left=0,
                        top=0,
                        width=builder.prs.slide_width,
                        height=builder.prs.slide_height
                    )
                except Exception as e:
                    logger.error(f"Failed to add background: {e}")
            else:
                # 回退到原圖
                logger.info(f"    使用原圖作為背景: {editable_img.image_path}")
                try:
                    slide.shapes.add_picture(
                        editable_img.image_path,
                        left=0,
                        top=0,
                        width=builder.prs.slide_width,
                        height=builder.prs.slide_height
                    )
                except Exception as e:
                    logger.error(f"Failed to add background: {e}")
            
            # 新增所有元素（遞迴地）
            # 計算縮放比例：將原始圖片座標對映到統一的幻燈片座標
            # 背景圖已經縮放到幻燈片尺寸，所以元素座標也需要相應縮放
            scale_x = slide_width_pixels / editable_img.width
            scale_y = slide_height_pixels / editable_img.height
            logger.info(f"    元素數量: {len(editable_img.elements)}, 圖片尺寸: {editable_img.width}x{editable_img.height}, "
                       f"幻燈片尺寸: {slide_width_pixels}x{slide_height_pixels}, 縮放比例: {scale_x:.3f}x{scale_y:.3f}")
            
            ExportService._add_editable_elements_to_slide(
                builder=builder,
                slide=slide,
                elements=editable_img.elements,
                scale_x=scale_x,
                scale_y=scale_y,
                depth=0,
                text_styles_cache=text_styles_cache,  # 使用預提取的樣式快取
                warnings=warnings  # 收集警告
            )
            
            logger.info(f"    ✓ 第 {page_idx + 1} 頁完成，新增了 {len(editable_img.elements)} 個元素")
        
        # 5. 儲存或返回位元組流
        report_progress("儲存檔案", "正在儲存PPTX檔案...", 95)
        if output_file:
            builder.save(output_file)
            report_progress("完成", f"✓ 可編輯PPTX已儲存", 100)
            logger.info(f"✓ 可編輯PPTX已儲存: {output_file}")
            
            # 輸出警告摘要
            if warnings.has_warnings():
                logger.warning(f"匯出完成，但有 {len(warnings.to_summary())} 條警告")
            
            return None, warnings
        else:
            pptx_bytes = builder.to_bytes()
            report_progress("完成", f"✓ 可編輯PPTX已生成", 100)
            logger.info(f"✓ 可編輯PPTX已生成（{len(pptx_bytes)} 位元組）")
            
            # 輸出警告摘要
            if warnings.has_warnings():
                logger.warning(f"匯出完成，但有 {len(warnings.to_summary())} 條警告")
            
            return pptx_bytes, warnings
    
    @staticmethod
    def _add_editable_elements_to_slide(
        builder,
        slide,
        elements: List,  # List[EditableElement]
        scale_x: float = 1.0,
        scale_y: float = 1.0,
        depth: int = 0,
        text_styles_cache: Dict[str, Any] = None,  # 預提取的文字樣式快取，key為element_id
        warnings: 'ExportWarnings' = None  # 警告收集器
    ):
        """
        遞迴地將EditableElement新增到幻燈片
        
        Args:
            builder: PPTXBuilder例項
            slide: 幻燈片物件
            elements: EditableElement列表
            scale_x: X軸縮放因子
            scale_y: Y軸縮放因子
            depth: 當前遞迴深度
            text_styles_cache: 預提取的文字樣式快取（可選），由 _batch_extract_text_styles 生成
        
        Note:
            elem.image_path 現在是絕對路徑，無需額外的目錄引數
        """
        if text_styles_cache is None:
            text_styles_cache = {}
        
        for elem in elements:
            elem_type = elem.element_type
            
            # 根據深度決定使用區域性座標還是全域性座標
            # depth=0: 頂層元素，使用區域性座標（bbox）
            # depth>0: 子元素，需要使用全域性座標（bbox_global）
            if depth == 0:
                bbox = elem.bbox  # 頂層元素使用區域性座標
            else:
                bbox = elem.bbox_global if hasattr(elem, 'bbox_global') and elem.bbox_global else elem.bbox
            
            # 轉換BBox物件為列表並應用縮放
            bbox_list = [
                int(bbox.x0 * scale_x),
                int(bbox.y0 * scale_y),
                int(bbox.x1 * scale_x),
                int(bbox.y1 * scale_y)
            ]
            
            logger.info(f"{'  ' * depth}  新增元素: type={elem_type}, bbox={bbox_list}, content={elem.content[:30] if elem.content else None}, image_path={elem.image_path}, 使用{'全域性' if depth > 0 else '區域性'}座標")
            
            # 根據型別新增元素（參考原實現的_add_mineru_text_to_slide和_add_mineru_image_to_slide）
            if elem_type in ['text', 'title', 'list', 'paragraph', 'header', 'footer', 'heading', 'table_caption', 'image_caption']:
                # 新增文字（參考_add_mineru_text_to_slide）
                if elem.content:
                    text = elem.content.strip()
                    if text:
                        try:
                            # 確定文字級別
                            level = 'title' if elem_type in ['title', 'heading'] else 'default'
                            
                            # 從快取獲取預提取的文字樣式
                            text_style = text_styles_cache.get(elem.element_id)
                            if text_style:
                                logger.debug(f"{'  ' * depth}  使用快取的文字樣式: color={text_style.font_color_rgb}, bold={text_style.is_bold}")
                            
                            builder.add_text_element(
                                slide=slide,
                                text=text,
                                bbox=bbox_list,
                                text_level=level,
                                text_style=text_style
                            )
                        except Exception as e:
                            logger.warning(f"新增文字元素失敗: {e}")
                            if warnings:
                                warnings.add_text_render_failed(text, str(e))
            
            elif elem_type == 'table_cell':
                # 新增表格單元格（帶邊框的文字框）
                if elem.content:
                    text = elem.content.strip()
                    if text:
                        try:
                            # 從快取獲取預提取的文字樣式
                            text_style = text_styles_cache.get(elem.element_id)
                            
                            # 表格單元格已經在上面統一處理了bbox_global和縮放
                            # 直接使用bbox_list即可
                            builder.add_text_element(
                                slide=slide,
                                text=text,
                                bbox=bbox_list,
                                text_level=None,
                                align='center',
                                text_style=text_style
                            )
                            
                        except Exception as e:
                            logger.warning(f"新增單元格失敗: {e}")
                            if warnings:
                                warnings.add_text_render_failed(text, str(e))
            
            elif elem_type == 'table':
                # 如果表格有子元素（單元格），使用inpainted背景 + 單元格
                if elem.children and elem.inpainted_background_path:
                    logger.info(f"{'  ' * depth}    表格有 {len(elem.children)} 個單元格，使用可編輯格式")
                    
                    # 先新增inpainted背景（乾淨的表格框架）
                    if os.path.exists(elem.inpainted_background_path):
                        try:
                            builder.add_image_element(
                                slide=slide,
                                image_path=elem.inpainted_background_path,
                                bbox=bbox_list
                            )
                        except Exception as e:
                            logger.error(f"Failed to add table background: {e}")
                    
                    # 遞迴新增單元格
                    ExportService._add_editable_elements_to_slide(
                        builder=builder,
                        slide=slide,
                        elements=elem.children,
                        scale_x=scale_x,
                        scale_y=scale_y,
                        depth=depth + 1,
                        text_styles_cache=text_styles_cache,
                        warnings=warnings
                    )
                else:
                    # 沒有子元素，新增整體表格圖片
                    # elem.image_path 現在是絕對路徑
                    if elem.image_path and os.path.exists(elem.image_path):
                        try:
                            builder.add_image_element(
                                slide=slide,
                                image_path=elem.image_path,
                                bbox=bbox_list
                            )
                        except Exception as e:
                            logger.error(f"Failed to add table image: {e}")
                    else:
                        logger.warning(f"Table image not found: {elem.image_path}")
                        builder.add_image_placeholder(slide, bbox_list)
            
            elif elem_type in ['image', 'figure', 'chart']:
                # 檢查是否應該使用遞迴渲染
                should_use_recursive_render = False
                
                if elem.children and elem.inpainted_background_path:
                    # 檢查是否有任意子元素佔據父元素絕大部分面積
                    parent_area = (bbox.x1 - bbox.x0) * (bbox.y1 - bbox.y0)
                    max_child_coverage_ratio = 0.85  # 閾值
                    has_dominant_child = False
                    
                    for child in elem.children:
                        if hasattr(child, 'bbox_global') and child.bbox_global:
                            child_bbox = child.bbox_global
                        else:
                            child_bbox = child.bbox
                        
                        child_area = child_bbox.area
                        coverage_ratio = child_area / parent_area if parent_area > 0 else 0
                        
                        if coverage_ratio > max_child_coverage_ratio:
                            logger.info(f"{'  ' * depth}    子元素 {child.element_id} 佔父元素面積 {coverage_ratio*100:.1f}% (>{max_child_coverage_ratio*100:.0f}%)，跳過遞迴渲染，直接使用原圖")
                            has_dominant_child = True
                            break
                    
                    should_use_recursive_render = not has_dominant_child
                
                # 如果有子元素且應該遞迴渲染
                if should_use_recursive_render:
                    logger.debug(f"{'  ' * depth}    元素有 {len(elem.children)} 個子元素，遞迴新增")
                    
                    # 先新增inpainted背景
                    if os.path.exists(elem.inpainted_background_path):
                        try:
                            builder.add_image_element(slide, elem.inpainted_background_path, bbox_list)
                        except Exception as e:
                            logger.error(f"Failed to add inpainted background: {e}")
                    
                    # 遞迴新增子元素
                    ExportService._add_editable_elements_to_slide(
                        builder=builder,
                        slide=slide,
                        elements=elem.children,
                        scale_x=scale_x,
                        scale_y=scale_y,
                        depth=depth + 1,
                        text_styles_cache=text_styles_cache,
                        warnings=warnings
                    )
                else:
                    # 沒有子元素或子元素佔比過大，直接新增原圖
                    # elem.image_path 現在是絕對路徑
                    if elem.image_path and os.path.exists(elem.image_path):
                        try:
                            builder.add_image_element(
                                slide=slide,
                                image_path=elem.image_path,
                                bbox=bbox_list
                            )
                        except Exception as e:
                            logger.error(f"Failed to add image: {e}")
                    else:
                        logger.warning(f"Image file not found: {elem.image_path}")
                        builder.add_image_placeholder(slide, bbox_list)
            
            else:
                # 其他型別
                logger.debug(f"{'  ' * depth}  跳過未知型別: {elem_type}")
    

