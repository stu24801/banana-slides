from typing import Dict
import json
from textwrap import dedent
from concurrent.futures import ThreadPoolExecutor, as_completed
from pathlib import Path
import re
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches
from gemini_genai import gen_image, gen_json_text, gen_text
def gen_outline(idea_prompt:str)->list[dict]:
    """generate outline of ppt, including optional parts and pages with title and points"""
    outline_prompt = dedent(f"""\
    You are a helpful assistant that generates an outline for a ppt.
    
    You can organize the content in two ways:
    
    1. Simple format (for short PPTs without major sections):
    [{{"title": "title1", "points": ["point1", "point2"]}}, {{"title": "title2", "points": ["point1", "point2"]}}]
    
    2. Part-based format (for longer PPTs with major sections):
    [
      {{
        "part": "Part 1: Introduction",
        "pages": [
          {{"title": "Welcome", "points": ["point1", "point2"]}},
          {{"title": "Overview", "points": ["point1", "point2"]}}
        ]
      }},
      {{
        "part": "Part 2: Main Content",
        "pages": [
          {{"title": "Topic 1", "points": ["point1", "point2"]}},
          {{"title": "Topic 2", "points": ["point1", "point2"]}}
        ]
      }}
    ]
    
    Choose the format that best fits the content. Use parts when the PPT has clear major sections.
    
    The user's request: {idea_prompt}. Now generate the outline, don't include any other text.
    使用全中文輸出。
    """)
    outline = gen_json_text(outline_prompt)
    outline = json.loads(outline)
    return outline
    

def flatten_outline(outline: list[dict]) -> list[dict]:
    """將可能包含part結構的outline扁平化為頁面列表"""
    pages = []
    for item in outline:
        if "part" in item and "pages" in item:
            # 這是一個part，展開其中的頁面
            for page in item["pages"]:
                # 為每個頁面新增part資訊
                page_with_part = page.copy()
                page_with_part["part"] = item["part"]
                pages.append(page_with_part)
        else:
            # 這是一個直接的頁面
            pages.append(item)
    return pages

def gen_desc(idea_prompt, outline: list[Dict])->list[Dict] :
    """generate description for each page, including title, full text content and more (並行生成)"""
    # 先將outline扁平化為頁面列表
    pages = flatten_outline(outline)
    
    # 為每個頁面準備生成任務
    def generate_page_desc(i, page_outline):
        part_info = f"\nThis page belongs to: {page_outline['part']}" if 'part' in page_outline else ""
        desc_prompt = dedent(f"""\
        we are generating the text desciption for each ppt page.
        the original user request is: \n{idea_prompt}\n
        We already have the entire ouline: \n{outline}\n{part_info}
        Now please generate the description for page {i}:
        {page_outline}
        The description includes page title, text to render(keep it concise).
        For example:
        頁面標題：原始社會：與自然共生
        頁面文字：
        - 狩獵採集文明： 人類活動規模小，對環境影響有限。
        - 依賴性強： 生活完全依賴於自然資源的直接供給，對自然規律敬畏。
        - 適應而非改造： 透過觀察和模仿學習自然，發展出適應當地環境的生存技能。
        - 影響特點： 區域性、短期、低強度，生態系統有充足的自我恢復能力。
        
        使用全中文輸出。
        """)
        page_desc = gen_text(desc_prompt)
        # 清理多餘的縮排
        page_desc = dedent(page_desc)
        return (i, page_desc)  # 返回索引和描述，以便排序
    
    # 使用執行緒池並行生成所有頁面的描述
    desc_dict = {}
    with ThreadPoolExecutor(max_workers=5) as executor:
        # 提交所有任務
        futures = [executor.submit(generate_page_desc, i, page_outline) 
                   for i, page_outline in enumerate(pages, 1)]
        
        # 收集結果
        for future in as_completed(futures):
            i, page_desc = future.result()
            desc_dict[i] = page_desc
            print(f"✓ 頁面 {i}/{len(pages)} 描述生成完成")
    
    # 按照原始順序返回結果
    desc = [desc_dict[i] for i in sorted(desc_dict.keys())]
    return desc

def gen_outline_text(outline: list[Dict]) -> str:
    """將outline轉換為文字格式，用於提示詞"""
    text_parts = []
    for i, item in enumerate(outline, 1):
        if "part" in item and "pages" in item:
            text_parts.append(f"{i}. {item['part']}")
        else:
            text_parts.append(f"{i}. {item.get('title', 'Untitled')}")
    result = "\n".join(text_parts)
    # 清理多餘的縮排
    return dedent(result)

def gen_prompts(outline: list[Dict], desc: list[str]) -> list[str]:
    """為每頁描述生成圖片提示詞"""
    pages = flatten_outline(outline)
    outline_text = gen_outline_text(outline)
    
    prompts = []
    for i, (page, page_desc) in enumerate(zip(pages, desc), 1):
        # 確定當前所屬章節
        if 'part' in page:
            current_section = page['part']
        else:
            current_section = f"{page.get('title', 'Untitled')}"
        
        # 構建提示詞，參考generate-example.py的格式
        prompt = dedent(f"""\
        利用專業平面設計知識，根據參考圖片的色彩與風格生成一頁設計風格相同的ppt頁面，作為整個ppt的其中一頁，內容是:
        {page_desc}
        
        整個ppt的大綱為：
        {outline_text}
        
        當前位於章節：{current_section}
        
        要求文字清晰銳利，畫面為4k解析度 16:9比例.畫面風格與配色保持嚴格一致。ppt使用全中文。
        """)
        print(f"\n-----\n prompt{i}:\n {prompt}\n-----\n")
        prompts.append(prompt)
    
    return prompts

def gen_images_parallel(prompts: list[str], ref_image: str, output_dir: str = "output") -> list[str]:
    """並行生成所有PPT頁面圖片"""
    # 建立輸出目錄
    output_path = Path(output_dir)
    output_path.mkdir(exist_ok=True)
    
    def generate_single_image(i, prompt):
        """生成單張圖片"""
        try:
            print(f"🎨 開始生成頁面 {i}/{len(prompts)} 的圖片...")
            image = gen_image(prompt, ref_image)
            if image:
                output_file = output_path / f"slide_{i:02d}.png"
                image.save(str(output_file))
                print(f"✓ 頁面 {i}/{len(prompts)} 圖片生成完成: {output_file}")
                return (i, str(output_file))
            else:
                print(f"✗ 頁面 {i}/{len(prompts)} 圖片生成失敗")
                return (i, None)
        except Exception as e:
            print(f"✗ 頁面 {i}/{len(prompts)} 生成出錯: {e}")
            return (i, None)
    
    # 使用執行緒池並行生成所有圖片
    image_files = {}
    with ThreadPoolExecutor(max_workers=8) as executor:  # 限制併發數為3避免API限流
        # 提交所有任務
        futures = [executor.submit(generate_single_image, i, prompt) 
                   for i, prompt in enumerate(prompts, 1)]
        
        # 收集結果
        for future in as_completed(futures):
            i, image_file = future.result()
            image_files[i] = image_file
    
    # 按照原始順序返回結果
    return [image_files[i] for i in sorted(image_files.keys())]

def create_pptx_from_images(input_dir: str = "output", output_file: str = "presentation.pptx"):
    """
    將指定目錄下的slide_XX.png圖片按順序組合成PPTX檔案
    
    Args:
        input_dir: 輸入圖片所在目錄
        output_file: 輸出的PPTX檔名
    """
    input_path = Path(input_dir)
    slide_files = list(input_path.glob("slide_*.png"))
    
    def extract_number(filename):
        match = re.search(r'slide_(\d+)', filename.stem)
        return int(match.group(1)) if match else 0
    
    slide_files.sort(key=extract_number)
    
    print(f"\n📁 找到 {len(slide_files)} 張幻燈片圖片")
    print(f"📝 開始建立 PPTX 檔案...")
    
    # 建立簡報
    prs = Presentation()
    
    # 設定幻燈片尺寸為16:9 (寬10英寸，高5.625英寸)
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    
    # 為每張圖片建立一頁幻燈片
    for i, image_file in enumerate(slide_files, 1):
        print(f"  ✓ 新增第 {i} 頁: {image_file.name}")
        
        # 新增空白幻燈片佈局（完全空白，沒有任何佔位符）
        blank_slide_layout = prs.slide_layouts[6]  # 佈局6通常是空白布局
        slide = prs.slides.add_slide(blank_slide_layout)
        
        # 將圖片新增到幻燈片，填充整個頁面
        # 左上角位置(0,0)，尺寸為幻燈片的完整寬高
        slide.shapes.add_picture(
            str(image_file),
            left=0,
            top=0,
            width=prs.slide_width,
            height=prs.slide_height
        )
    
    # 儲存PPTX檔案
    prs.save(output_file)
    
    print(f"\n✅ 成功建立 PPTX 檔案: {output_file}")
    print(f"📊 總共 {len(slide_files)} 頁幻燈片")
    return True

def gen_ppt(idea_prompt, ref_image):
    # 建立帶時間戳的輸出目錄
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    output_dir = f"output_{timestamp}"
    pptx_filename = f"presentation_{timestamp}.pptx"
    
    print(f"📂 本次執行輸出目錄: {output_dir}")
    print(f"📄 PPTX檔名: {pptx_filename}\n")
    
    outline = gen_outline(idea_prompt)
    
    # 顯示原始outline結構（可能包含parts）
    print("PPT Outline:")
    for item in outline:
        if "part" in item and "pages" in item:
            print(f"\n【{item['part']}】")
            for j, page in enumerate(item["pages"], 1):
                print(f"  Page {j}: {page.get('title', 'Untitled')}")
                print(f"    Points: {page.get('points', [])}")
        else:
            print(f"\nPage: {item.get('title', 'Untitled')}")
            print(f"  Points: {item.get('points', [])}")
    
    # 生成詳細描述
    desc = gen_desc(idea_prompt, outline)
    
    # 顯示每頁描述
    pages = flatten_outline(outline)
    for i, (page, page_desc) in enumerate(zip(pages, desc), 1):
        part_tag = f"[{page['part']}] " if 'part' in page else ""
        print(f"-----\nPage {i} {part_tag}- {page.get('title', 'Untitled')}\n-----")
        print(f"{page_desc}\n")
    
    # 生成圖片提示詞
    print("開始生成圖片提示詞...")
    prompts = gen_prompts(outline, desc)
    print(f"✓ 已生成 {len(prompts)} 個頁面的提示詞\n")
    
    # 並行生成所有頁面圖片（使用帶時間戳的目錄）
    print("開始並行生成PPT頁面圖片...")
    image_files = gen_images_parallel(prompts, ref_image, output_dir)
    
    # 顯示結果彙總
    print("PPT圖片生成完成！")
    successful = [f for f in image_files if f is not None]
    print(f"✓ 成功生成 {len(successful)}/{len(image_files)} 張圖片")
    for i, image_file in enumerate(image_files, 1):
        if image_file:
            print(f"  頁面 {i}: {image_file}")
        else:
            print(f"  頁面 {i}: 生成失敗")
    
    # 將所有圖片組合成PPTX檔案
    if successful:
        print("正在生成最終的PPTX檔案...")
        create_pptx_from_images(output_dir, pptx_filename)
    
    return image_files
    
    

if __name__ == "__main__":
    idea_prompt="生成一張關於人類活動對生態環境影響的ppt.只要3頁。"
    ref_image="template_g.png"
    gen_ppt(idea_prompt, ref_image)